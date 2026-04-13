"""Document parsers — PDF, PPTX, Markdown, plain text."""

from __future__ import annotations

import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any


@dataclass
class DocumentChunk:
    """A single chunk (page, slide, section, paragraph) of a parsed document."""

    doc_id: str
    chunk_index: int
    chunk_label: str
    text: str
    metadata: dict[str, Any] = field(default_factory=dict)


def parse_pdf(path: str, doc_id: str) -> list[DocumentChunk]:
    """Parse a PDF into one chunk per page using pymupdf."""
    try:
        import fitz  # pymupdf
    except ImportError:
        raise ImportError("pymupdf is required for PDF support: pip install pymupdf")

    chunks: list[DocumentChunk] = []
    with fitz.open(path) as doc:
        for i, page in enumerate(doc):
            text = page.get_text().strip()
            if not text:
                continue
            chunks.append(
                DocumentChunk(
                    doc_id=doc_id,
                    chunk_index=i,
                    chunk_label=f"page {i + 1}",
                    text=text,
                    metadata={
                        "doc_id": doc_id,
                        "chunk_index": i,
                        "chunk_label": f"page {i + 1}",
                        "file_type": "pdf",
                        "file_path": path,
                        "page_number": i + 1,
                    },
                )
            )
    return chunks


def parse_pptx(path: str, doc_id: str) -> list[DocumentChunk]:
    """Parse a PPTX into one chunk per slide."""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("python-pptx is required for PPTX support: pip install python-pptx")

    chunks: list[DocumentChunk] = []
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = paragraph.text.strip()
                    if line:
                        texts.append(line)
        text = "\n".join(texts)
        if not text:
            text = "[no text content]"
        chunks.append(
            DocumentChunk(
                doc_id=doc_id,
                chunk_index=i,
                chunk_label=f"slide {i + 1}",
                text=text,
                metadata={
                    "doc_id": doc_id,
                    "chunk_index": i,
                    "chunk_label": f"slide {i + 1}",
                    "file_type": "pptx",
                    "file_path": path,
                    "slide_number": i + 1,
                },
            )
        )
    return chunks


def parse_markdown(path: str, doc_id: str) -> list[DocumentChunk]:
    """Parse a Markdown file into one chunk per heading section."""
    content = Path(path).read_text(encoding="utf-8")
    heading_pattern = re.compile(r"^(#{1,6})\s+(.+)", re.MULTILINE)

    splits: list[tuple[int, str]] = []
    for m in heading_pattern.finditer(content):
        splits.append((m.start(), m.group(2).strip()))

    if not splits:
        # No headings — single chunk
        text = content.strip()
        if not text:
            return []
        return [
            DocumentChunk(
                doc_id=doc_id,
                chunk_index=0,
                chunk_label="full document",
                text=text,
                metadata={
                    "doc_id": doc_id,
                    "chunk_index": 0,
                    "chunk_label": "full document",
                    "file_type": "md",
                    "file_path": path,
                },
            )
        ]

    chunks: list[DocumentChunk] = []

    # Content before the first heading
    preamble = content[: splits[0][0]].strip()
    if preamble:
        chunks.append(
            DocumentChunk(
                doc_id=doc_id,
                chunk_index=0,
                chunk_label="preamble",
                text=preamble,
                metadata={
                    "doc_id": doc_id,
                    "chunk_index": 0,
                    "chunk_label": "preamble",
                    "file_type": "md",
                    "file_path": path,
                },
            )
        )

    for idx, (start, heading) in enumerate(splits):
        end = splits[idx + 1][0] if idx + 1 < len(splits) else len(content)
        text = content[start:end].strip()
        if not text:
            continue
        ci = len(chunks)
        label = f"section: {heading}"
        chunks.append(
            DocumentChunk(
                doc_id=doc_id,
                chunk_index=ci,
                chunk_label=label,
                text=text,
                metadata={
                    "doc_id": doc_id,
                    "chunk_index": ci,
                    "chunk_label": label,
                    "file_type": "md",
                    "file_path": path,
                    "section_heading": heading,
                },
            )
        )

    return chunks


def parse_epub(path: str, doc_id: str) -> list[DocumentChunk]:
    """Parse an EPUB into one chunk per chapter/spine item."""
    try:
        import ebooklib
        from ebooklib import epub
    except ImportError:
        raise ImportError("ebooklib is required for EPUB support: pip install ebooklib beautifulsoup4")

    try:
        from bs4 import BeautifulSoup
    except ImportError:
        raise ImportError("beautifulsoup4 is required for EPUB support: pip install beautifulsoup4")

    book = epub.read_epub(path)

    # Build a map of title/heading from the TOC for each spine item href.
    toc_titles: dict[str, str] = {}

    def _walk_toc(entries: Any) -> None:
        for entry in entries:
            if isinstance(entry, tuple):
                section, children = entry
                if hasattr(section, "href") and hasattr(section, "title"):
                    href = section.href.split("#")[0]
                    toc_titles.setdefault(href, section.title)
                _walk_toc(children)
            elif hasattr(entry, "href") and hasattr(entry, "title"):
                href = entry.href.split("#")[0]
                toc_titles.setdefault(href, entry.title)

    _walk_toc(book.toc)

    chunks: list[DocumentChunk] = []
    for item in book.get_items():
        if item.get_type() != ebooklib.ITEM_DOCUMENT:
            continue

        soup = BeautifulSoup(item.get_body_content(), "html.parser")
        # Remove script/style
        for tag in soup(["script", "style"]):
            tag.decompose()
        text = soup.get_text(separator="\n", strip=True)
        # Collapse excessive blank lines
        text = re.sub(r"\n{3,}", "\n\n", text).strip()

        if not text or len(text) < 20:
            continue

        href = item.get_name()
        title = toc_titles.get(href)
        if not title:
            # Fall back to first heading in the chapter
            heading = soup.find(["h1", "h2", "h3"])
            title = heading.get_text(strip=True) if heading else href

        ci = len(chunks)
        label = f"chapter: {title}"
        chunks.append(
            DocumentChunk(
                doc_id=doc_id,
                chunk_index=ci,
                chunk_label=label,
                text=text,
                metadata={
                    "doc_id": doc_id,
                    "chunk_index": ci,
                    "chunk_label": label,
                    "file_type": "epub",
                    "file_path": path,
                    "chapter_title": title,
                    "chapter_href": href,
                },
            )
        )

    return chunks


def parse_text(path: str, doc_id: str) -> list[DocumentChunk]:
    """Parse a plain text file into one chunk per paragraph."""
    content = Path(path).read_text(encoding="utf-8")
    paragraphs = re.split(r"\n\n+", content)

    chunks: list[DocumentChunk] = []
    for i, para in enumerate(paragraphs):
        text = para.strip()
        if not text:
            continue
        label = f"paragraph {i + 1}"
        chunks.append(
            DocumentChunk(
                doc_id=doc_id,
                chunk_index=len(chunks),
                chunk_label=label,
                text=text,
                metadata={
                    "doc_id": doc_id,
                    "chunk_index": len(chunks),
                    "chunk_label": label,
                    "file_type": "txt",
                    "file_path": path,
                    "paragraph_number": i + 1,
                },
            )
        )
    return chunks


_EXTENSION_MAP: dict[str, str] = {
    ".pdf": "pdf",
    ".pptx": "pptx",
    ".ppt": "pptx",
    ".md": "md",
    ".markdown": "md",
    ".epub": "epub",
    ".txt": "txt",
    ".text": "txt",
    ".log": "txt",
    ".csv": "txt",
    ".json": "txt",
    ".yaml": "txt",
    ".yml": "txt",
    ".toml": "txt",
}

_PARSERS = {
    "pdf": parse_pdf,
    "pptx": parse_pptx,
    "md": parse_markdown,
    "epub": parse_epub,
    "txt": parse_text,
}


def parse_document(path: str, existing_ids: set[str] | None = None) -> list[DocumentChunk]:
    """Parse a document, inferring type from extension.

    Args:
        path: Path to the document file.
        existing_ids: Already-used doc_ids to avoid collisions.
    """
    p = Path(path).resolve()
    if not p.exists():
        raise FileNotFoundError(f"Document not found: {path}")

    ext = p.suffix.lower()
    file_type = _EXTENSION_MAP.get(ext)
    if file_type is None:
        raise ValueError(f"Unsupported file type: {ext} (supported: {', '.join(sorted(_EXTENSION_MAP))})")

    # Generate unique doc_id
    base_id = re.sub(r"[^a-zA-Z0-9_-]", "_", p.stem)
    doc_id = base_id
    used = existing_ids or set()
    counter = 2
    while doc_id in used:
        doc_id = f"{base_id}-{counter}"
        counter += 1

    parser = _PARSERS[file_type]
    return parser(str(p), doc_id)
